from pptx import Presentation
import textwrap

def shorten_text(text, max_length):
    """
    Shorten the text to the max_length, adding an ellipsis if the text is shortened.
    """
    return textwrap.shorten(text, width=max_length, placeholder="...")

def process_slide(slide, max_length):
    """
    Process each slide, shortening the text in each text box if it's too long.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            if len(paragraph.text) > max_length:
                paragraph.text = shorten_text(paragraph.text, max_length)

def process_presentation(input_file, output_file, max_length):
    """
    Process each slide in the presentation.
    """
    ppt = Presentation(input_file)
    for slide in ppt.slides:
        process_slide(slide, max_length)
    ppt.save(output_file)

# Example usage
input_file = 'path/to/input/presentation.pptx'
output_file = 'path/to/output/presentation.pptx'
max_length = 50  # Maximum length of text in characters
process_presentation(input_file, output_file, max_length)
