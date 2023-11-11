import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# Function to extract features from a single presentation
def extract_features(ppt_path):
    prs = Presentation(ppt_path)
    features = {
        'font_size': [],  # Placeholder for font sizes
        'colors': [],  # Placeholder for colors
        'font_family': []  # Placeholder for font families
    }

    # Example: Iterate through slides and extract features
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            features['font_size'].append(run.font.size)
                        if run.font.name:
                            features['font_family'].append(run.font.name)

            # Add logic for extracting image ratios, margins, alignments, colors, etc.

    return features

# Load scores from the text file
# Replace with actual logic to load scores
scores = {}  

# Array of PowerPoint file paths
ppt_paths = ['./p'+str(x)+'.pptx' for x in range(1,4)]

# Create a DataFrame to hold all features and scores
data = []

# Iterate through each PowerPoint file path in the array and extract features
for ppt_path in ppt_paths:
    features = extract_features(ppt_path)
    features['score'] = scores.get(ppt_path, None)  # Get the score for the presentation
    data.append(features)

# Convert to DataFrame
df = pd.DataFrame(data)

# Save the DataFrame to a CSV file
df.to_csv('presentation_data.csv', index=False)
