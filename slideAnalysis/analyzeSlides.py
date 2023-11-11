'''

{
slides: [
 {id:1
  message:"the summary of slide 1",
  image:"the summary of image contents in slide 1"
 },
 {id:2
  message:"the summary of slide 2",
  image:"the summary of image contents in slide 2"
 },
]
}

'''

from pptx import Presentation
import os, json

filename = 'capsher.pptx'
foldername = filename[:-5] + '_data'
if not os.path.exists(foldername):
    os.mkdir(foldername)
presentation = Presentation(filename)

info = []
id = 0



for slide in presentation.slides:
    
    id += 1
    title = ''
    textBoxes = []
    for shape in slide.shapes:
        if shape.shape_type == 1:
            title = shape.text
        elif shape.has_text_frame:
            textBoxes.append(shape.text)
    text = {'title' : title, 'textboxes' : textBoxes}

    image_info = []
    image_number = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:
            image_number += 1
            image = shape.image
            image_filename = f'slide{id}-image{image_number}.png'
            image_info.append(image_filename)
            image_path = os.path.join(foldername, image_filename)
            with open(image_path, 'wb') as img_file:
                img_file.write(image.blob)
    info.append({'id' : id, 'text' : text, 'image' : image_info})

json_path = os.path.join(foldername, '.json')
# clear the file
with open(json_path, 'w'):
    pass
for thing in info:
    print(thing)
    with open(json_path,'a') as json_file:
        json.dump(thing, json_file)
        json_file.write('\n')
