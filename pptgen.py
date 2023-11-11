from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.picture import Picture

def replace_image(ppt_path, slide_count, image_shape_id, new_image_path):
    prs = Presentation(ppt_path)

    # Access the specific slide
    for i in range(slide_count):
        slide=prs.slides[i]
        cords = []
        for i, shape in enumerate(slide.shapes):
            if isinstance(shape, Picture):
                print("update")
                old_image = shape
                cords.append({'left':old_image.left, 'top':old_image.top, 'width':old_image.width, 'height':old_image.height}) 
                slide.shapes._spTree.remove(old_image._element)
        for cord in cords:
            slide.shapes.add_picture(new_image_path, cord['left'], cord['top'], width=cord['width'], height=cord['height'])
    prs.save('updated_presentation.pptx')

# Replace parameters with your specific details
ppt_path = 'p3.pptx'  # Path to your existing PowerPoint file
slide_count = 40  # Index of the slide containing the image (0 for the first slide)
image_shape_id = 0  # ID of the shape to be replaced (you need to find this ID)
new_image_path = 'wok.jpg'  # Path to the new image you want to insert

replace_image(ppt_path, slide_count, image_shape_id, new_image_path)
